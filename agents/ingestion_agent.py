"""
ingestion_agent.py
------------------
Document Ingestion Agent — Phase 1 of the AI automation pipeline.

Responsibilities:
  1. Accept any document file (PDF, DOCX, XLSX, PPTX, TXT, MD)
  2. Extract all readable text and table content
  3. Use LLM to structure the content into a rich feature_description
     and detect which SDP modules are likely involved
  4. Populate AgentState with feature_description + document_metadata
     so the downstream Planner Agent can work with zero manual input

If no source_document is present in state, this agent is a no-op —
the existing pipeline flow (text-based feature_description) is unchanged.

Supported formats:
  .pdf   — PyMuPDF (page text + embedded tables)
  .docx  — python-docx (paragraphs + tables)
  .xlsx  — openpyxl (all sheets, rows as scenario lines)
  .pptx  — python-pptx (slide text boxes + tables)
  .txt   — plain text read
  .md    — plain text read
"""

import json
import os
from pathlib import Path
from typing import Any

from langchain_core.messages import SystemMessage, HumanMessage

from agents.state import AgentState
from agents.llm_factory import get_llm

# Maximum characters fed to LLM (avoid token overflow for very large docs)
MAX_CHARS = 12_000


# ─────────────────────────────────────────────────────────────────────────────
# Text extractors
# ─────────────────────────────────────────────────────────────────────────────

def _extract_pdf(path: str) -> str:
    """Extract text from every page of a PDF using PyMuPDF."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text("text").strip()
            if text:
                pages.append(f"[Page {i + 1}]\n{text}")
        return "\n\n".join(pages)
    except ImportError:
        raise RuntimeError("PyMuPDF not installed. Run: pip install PyMuPDF")


def _extract_docx(path: str) -> str:
    """Extract paragraphs and tables from a DOCX file."""
    try:
        from docx import Document
        doc = Document(path)
        parts = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append("\n".join(rows))

        return "\n\n".join(parts)
    except ImportError:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")


def _extract_xlsx(path: str) -> str:
    """Extract all cell values from every sheet of an XLSX file."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip() not in ("", "None")]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        return "\n\n".join(parts)
    except ImportError:
        raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")


def _extract_pptx(path: str) -> str:
    """Extract text from all slides of a PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))
            if texts:
                slides.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")


def _extract_text(path: str) -> str:
    """Read a plain text or markdown file."""
    with open(path, encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_document_text(file_path: str) -> str:
    """
    Dispatch to the correct extractor based on file extension.
    Returns raw extracted text (may be large).
    """
    ext = Path(file_path).suffix.lower()
    extractors = {
        ".pdf":  _extract_pdf,
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".xls":  _extract_xlsx,
        ".pptx": _extract_pptx,
        ".txt":  _extract_text,
        ".md":   _extract_text,
        ".rst":  _extract_text,
    }
    fn = extractors.get(ext)
    if fn is None:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported: {', '.join(extractors.keys())}"
        )
    return fn(file_path)


# ─────────────────────────────────────────────────────────────────────────────
# LLM structuring prompt
# ─────────────────────────────────────────────────────────────────────────────

SYSTEM_PROMPT = """You are a senior QA architect for Zoho ServiceDesk Plus (SDP),
an enterprise ITSM platform. You specialise in reading feature specifications,
user stories, BA documents, test plans, and requirement sheets, then converting
them into structured test scenario descriptions that can be automated.

SDP module areas include:
  requests (incident/service requests), problems, changes, releases, assets,
  cmdb, solutions, projects, tasks, admin (automation, notification rules, SLA,
  workflows), purchase, contracts, reports, dashboard, self-service portal.

Your job:
1. Read the provided document text
2. Identify the SDP module(s) / feature(s) being described
3. Extract every testable scenario, acceptance criterion, or edge case mentioned
4. Produce a detailed, well-structured feature_description text that a downstream
   QA Planner Agent can use to generate complete test cases
5. Suggest module paths using the format: "requests/request", "solutions/solution",
   "admin/automation/notificationrules", etc.

Respond ONLY with a valid JSON object in this exact format:
{
  "document_title": "Short title extracted or inferred from the document",
  "document_type": "feature_spec | user_story | test_plan | ba_document | requirement_sheet | unknown",
  "feature_description": "A detailed multi-paragraph description of the feature and all its testable behaviours. Include: what the feature does, what inputs/configurations it accepts, what the expected outcomes are, edge cases mentioned, role-based access differences, and any validation rules. This should be rich enough for a QA planner to generate 5-20 test scenarios from it.",
  "detected_modules": ["module/path1", "module/path2"],
  "extracted_scenarios": [
    {
      "title": "Short scenario title",
      "description": "Detailed description of this specific scenario",
      "acceptance_criteria": ["Criterion 1", "Criterion 2"],
      "type": "CREATE | READ | UPDATE | DELETE | VALIDATE | ROLE_BASED | NEGATIVE | EDGE_CASE"
    }
  ],
  "confidence": "high | medium | low",
  "notes": "Any caveats, ambiguities, or missing information in the document"
}"""


# ─────────────────────────────────────────────────────────────────────────────
# Test Case Sheet extraction prompt
# Used when the uploaded document IS already a test case register (e.g. Excel QA
# test case sheet) rather than a feature spec.  Bypasses Planner + Coverage.
# ─────────────────────────────────────────────────────────────────────────────

TESTCASE_SYSTEM_PROMPT = """You are a senior QA engineer for Zoho ServiceDesk Plus (SDP),
an enterprise ITSM platform. You are given the raw contents of a QA test case document
(typically an Excel sheet, CSV, or tabular test register).

Your job is to extract each individual test case EXACTLY as written — do NOT invent or
interpret. Map every test case to its SDP module path and produce a camelCase Java method
name from the test case title.

SDP module path conventions:
  "requests/request"            — Incident / Service requests
  "problems/problem"            — Problem management
  "changes/change"              — Change management
  "releases/release"            — Release management
  "assets/asset"                — Asset management
  "solutions/solution"          — Solutions / Knowledge base
  "projects/project"            — Project management
  "tasks/task"                  — Task management
  "admin/automation/workflows"  — Workflow automation
  "admin/automation/notificationrules" — Notification rules
  "admin/sla"                   — SLA management
  "admin/zia"                   — ZIA (AI) settings
  "reports/report"              — Reports

Rules:
1. Extract EVERY test case row — do not skip or merge any.
2. method_name must be unique camelCase, derived from the title (e.g. "verifyCreateIncidentWithHighPriority").
3. group should be "NoPreprocess" unless the test explicitly needs a prerequisite entity
   (in that case use "CREATE_PREREQUISITE").
4. test_steps should be a list of individual action strings from the test steps column.
5. expected_result must be the verbatim expected outcome text.
6. data_requirements: list field names that the test needs (e.g. ["priority", "category", "subject"]).

Respond ONLY with a valid JSON object in this exact format (no markdown fences):
{
  "document_title": "Short title of the test case sheet",
  "scenarios": [
    {
      "module_path": "requests/request",
      "title": "Exact test case title from the document",
      "method_name": "camelCaseMethodName",
      "description": "One-paragraph summary of what this test validates",
      "test_steps": ["Navigate to ...", "Click ...", "Fill in ..."],
      "expected_result": "Verbatim expected outcome",
      "group": "NoPreprocess",
      "data_requirements": ["subject", "priority"],
      "notes": "Any extra info or caveats for this test case"
    }
  ]
}"""

# Document types that indicate an already-written test case register
_TESTCASE_DOC_TYPES = {"test_plan", "testcase_sheet"}


# ─────────────────────────────────────────────────────────────────────────────
# IngestionAgent
# ─────────────────────────────────────────────────────────────────────────────

class IngestionAgent:
    """
    LangGraph node: Document → structured feature_description.

    Usage as a standalone pre-processor:
        agent = IngestionAgent()
        result = agent.process_file("path/to/feature.pdf")
        # result['feature_description'] → feed to run_pipeline(feature_description=...)

    Usage as a LangGraph node:
        State must include 'source_document' (file path string).
        Agent populates: feature_description, target_modules, document_metadata.
    """

    def __init__(self, llm: Any = None, base_dir: str = None):
        self.base = Path(base_dir) if base_dir else Path(__file__).resolve().parents[1]
        self.llm = llm or get_llm(temperature=0.1)

    def _call_llm(self, raw_text: str) -> dict:
        """Send extracted text to LLM for structuring. Returns parsed JSON dict."""
        # Truncate if necessary
        truncated = raw_text[:MAX_CHARS]
        if len(raw_text) > MAX_CHARS:
            truncated += f"\n\n[... document truncated at {MAX_CHARS} chars ...]"

        response = self.llm.invoke([
            SystemMessage(content=SYSTEM_PROMPT),
            HumanMessage(content=f"Document content:\n\n{truncated}"),
        ])

        raw = response.content.strip()
        # Strip markdown fences if model added them
        if raw.startswith("```"):
            raw = raw.split("\n", 1)[1]
            raw = raw.rsplit("```", 1)[0]

        return json.loads(raw)

    def _call_llm_testcases(self, raw_text: str) -> dict:
        """Use the testcase-specific prompt to extract structured proposed_scenarios."""
        truncated = raw_text[:MAX_CHARS]
        if len(raw_text) > MAX_CHARS:
            truncated += f"\n\n[... document truncated at {MAX_CHARS} chars ...]"

        response = self.llm.invoke([
            SystemMessage(content=TESTCASE_SYSTEM_PROMPT),
            HumanMessage(content=f"Test case document contents:\n\n{truncated}"),
        ])

        raw = response.content.strip()
        if raw.startswith("```"):
            raw = raw.split("\n", 1)[1]
            raw = raw.rsplit("```", 1)[0]

        return json.loads(raw)

    def process_file_testcases(self, file_path: str) -> dict:
        """
        Process a test case register document (Excel/CSV/tabular) and extract
        each test case directly as a proposed_scenario dict.

        Returns:
          - proposed_scenarios : list[dict]  — ready for CoderAgent (bypasses Planner/Coverage)
          - document_metadata  : dict
        """
        file_path = str(file_path)
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document not found: {file_path}")

        print(f"[IngestionAgent] Extracting test cases from: {Path(file_path).name}")
        raw_text = extract_document_text(file_path)
        print(f"[IngestionAgent] Extracted {len(raw_text):,} chars. Parsing test cases...")

        structured = self._call_llm_testcases(raw_text)
        scenarios  = structured.get("scenarios", [])

        print(f"[IngestionAgent] ✅ Extracted {len(scenarios)} test case(s) from document.")

        return {
            "proposed_scenarios": scenarios,
            "document_metadata": {
                "source_file":     file_path,
                "document_title":  structured.get("document_title", Path(file_path).stem),
                "document_type":   "testcase_sheet",
                "confidence":      "high",
                "notes":           f"{len(scenarios)} test cases extracted directly — bypassing Planner/Coverage.",
                "scenario_count":  len(scenarios),
                "raw_text_length": len(raw_text),
            },
        }


        """
        Process a document file end-to-end.

        Returns a dict with:
          - feature_description: str  (ready for PlannerAgent)
          - target_modules: list[str] (detected module hints)
          - document_metadata: dict   (title, type, scenarios, confidence, notes)
          - raw_text_length: int
        """
        file_path = str(file_path)
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document not found: {file_path}")

        print(f"[IngestionAgent] Extracting text from: {Path(file_path).name}")
        raw_text = extract_document_text(file_path)
        print(f"[IngestionAgent] Extracted {len(raw_text):,} characters. Structuring with LLM...")

        structured = self._call_llm(raw_text)

        feature_description = structured.get("feature_description", "")
        # Enrich feature_description with extracted_scenarios for maximum Planner context
        scenarios = structured.get("extracted_scenarios", [])
        if scenarios:
            scenario_block = "\n\nExtracted Scenarios from Document:\n"
            for i, sc in enumerate(scenarios, 1):
                scenario_block += f"\n{i}. [{sc.get('type', 'TEST')}] {sc.get('title', '')}\n"
                scenario_block += f"   {sc.get('description', '')}\n"
                for ac in sc.get("acceptance_criteria", []):
                    scenario_block += f"   - {ac}\n"
            feature_description += scenario_block

        return {
            "feature_description": feature_description,
            "target_modules": structured.get("detected_modules", []),
            "document_metadata": {
                "source_file":      file_path,
                "document_title":   structured.get("document_title", Path(file_path).stem),
                "document_type":    structured.get("document_type", "unknown"),
                "confidence":       structured.get("confidence", "medium"),
                "notes":            structured.get("notes", ""),
                "scenario_count":   len(scenarios),
                "raw_text_length":  len(raw_text),
            },
        }

    def run(self, state: AgentState) -> AgentState:
        """
        LangGraph node function.

        If 'source_document' is set in state: extract + structure the document,
        then populate feature_description, target_modules, document_metadata.

        If 'source_document' is NOT set: pass through unchanged (no-op).
        This preserves backward compatibility with all existing pipeline calls.
        """
        source_doc = state.get("source_document", "")
        if not source_doc:
            # No document provided — existing text-based flow, nothing to do
            state["messages"] = [
                "[IngestionAgent] No source_document in state — skipping (text-based flow)"
            ]
            return state

        state["messages"] = [
            f"[IngestionAgent] Processing document: {Path(source_doc).name}"
        ]

        try:
            # ── Testcase-sheet mode (bypass Planner + Coverage) ────────────
            # Triggered when caller explicitly sets generation_mode="from_testcases"
            # OR when the file is an Excel/CSV-type spreadsheet that we suspect
            # already contains structured test cases (auto-detection).
            is_explicit_tc_mode = state.get("generation_mode", "") == "from_testcases"
            is_spreadsheet      = Path(source_doc).suffix.lower() in (".xlsx", ".xls", ".csv")

            if is_explicit_tc_mode or (is_spreadsheet and not state.get("feature_description", "").strip()):
                # Run the dedicated testcase extractor
                tc_result = self.process_file_testcases(source_doc)

                if tc_result.get("proposed_scenarios"):
                    state["proposed_scenarios"]  = tc_result["proposed_scenarios"]
                    state["generation_mode"]     = "from_testcases"
                    state["document_metadata"]   = tc_result["document_metadata"]
                    meta = tc_result["document_metadata"]
                    state["messages"] = [
                        f"[IngestionAgent] ✅ Test case sheet processed: "
                        f"'{meta['document_title']}' — "
                        f"{meta['scenario_count']} test case(s) extracted. "
                        f"Pipeline will bypass Planner + Coverage + Scout."
                    ]
                    return state
                else:
                    # LLM returned no scenarios — fall through to normal feature-doc flow
                    print("[IngestionAgent] ⚠️  Testcase extraction returned 0 scenarios — falling back to feature-doc flow.")
                    state["generation_mode"] = "new_feature"

            result = self.process_file(source_doc)

            # Only overwrite feature_description if it was empty
            if not state.get("feature_description", "").strip():
                state["feature_description"] = result["feature_description"]
            else:
                # Merge: prepend document content to any existing description
                state["feature_description"] = (
                    result["feature_description"]
                    + "\n\nAdditional context provided:\n"
                    + state["feature_description"]
                )

            # Merge detected module hints with any already specified
            existing_modules = state.get("target_modules", []) or []
            detected = result["target_modules"]
            merged = list(dict.fromkeys(existing_modules + detected))  # preserve order, dedupe
            state["target_modules"] = merged

            state["document_metadata"] = result["document_metadata"]

            state["messages"] = [
                f"[IngestionAgent] ✅ Document processed: "
                f"'{result['document_metadata']['document_title']}' "
                f"({result['document_metadata']['document_type']}) — "
                f"{result['document_metadata']['scenario_count']} scenarios extracted, "
                f"{len(merged)} module(s) detected "
                f"[confidence: {result['document_metadata']['confidence']}]"
            ]

        except Exception as e:
            state["errors"] = [
                f"[IngestionAgent] Error processing '{source_doc}': {e}"
            ]
            state["messages"] = [
                f"[IngestionAgent] ❌ Failed to process document: {e}"
            ]

        return state
